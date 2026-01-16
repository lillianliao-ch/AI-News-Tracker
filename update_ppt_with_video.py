#!/usr/bin/env python3
"""
更新蔡徐坤PPT，添加视频链接页面
由于python-pptx视频嵌入支持有限，我们创建一个专门的视频页面
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_gradient_slide(slide):
    """为幻灯片添加深色背景"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.33), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(25, 25, 35)
    background.line.fill.background()
    
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title_text(slide, text, top, font_size=44, bold=True, color=RGBColor(255, 255, 255)):
    """添加标题文本"""
    left = Inches(0.5)
    width = Inches(12.33)
    height = Inches(1)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "PingFang SC"
    return textbox

def add_bullet_points(slide, points, top, left=Inches(1), font_size=24):
    """添加要点列表"""
    width = Inches(11)
    height = Inches(4)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        run = p.add_run()
        run.text = f"• {point}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(230, 230, 230)
        run.font.name = "PingFang SC"
        p.space_after = Pt(12)
    return textbox

def add_emoji_text(slide, emoji, text, top, font_size=28):
    """添加带emoji的文本"""
    left = Inches(0.5)
    width = Inches(12.33)
    height = Inches(0.8)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{emoji} {text}"
    run.font.size = Pt(font_size)
    run.font.color.rgb = RGBColor(255, 215, 0)
    run.font.name = "PingFang SC"
    return textbox

def create_caixukun_ppt_with_video():
    """创建带视频的蔡徐坤介绍PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6]
    
    # ========== 第1页：封面 ==========
    slide1 = prs.slides.add_slide(blank_layout)
    create_gradient_slide(slide1)
    
    add_title_text(slide1, "🏀", Inches(1.5), font_size=72, bold=False)
    add_title_text(slide1, "蔡徐坤", Inches(2.5), font_size=60)
    add_title_text(slide1, "「一个被唱跳耽误的篮球运动员」", Inches(3.5), font_size=28, bold=False, color=RGBColor(200, 200, 200))
    add_title_text(slide1, "——  轻松了解这位顶流偶像  ——", Inches(4.5), font_size=20, bold=False, color=RGBColor(150, 150, 150))
    
    # ========== 第2页：人物简介 ==========
    slide2 = prs.slides.add_slide(blank_layout)
    create_gradient_slide(slide2)
    
    add_title_text(slide2, "👤 人物档案", Inches(0.5), font_size=40)
    
    bio_points = [
        "🎂 1998年8月2日出生于浙江温州（著名的「温州皮革厂」老乡）",
        "🎓 美国加州格雷斯兄弟高中毕业（留学海归，不简单！）",
        "🎤 身份：歌手 / 原创音乐制作人 / 演员 / ...... 篮球爱好者？",
        "📱 微博粉丝：3000万+（比很多小国家人口还多）",
        "🏆 2018年《偶像练习生》C位出道，NINE PERCENT队长",
    ]
    add_bullet_points(slide2, bio_points, Inches(1.5), font_size=26)
    
    add_emoji_text(slide2, "💡", "冷知识：差点成为TFBOYS成员，被妈妈送去美国读书了", Inches(6.3), font_size=20)
    
    # ========== 第3页：出道历程 ==========
    slide3 = prs.slides.add_slide(blank_layout)
    create_gradient_slide(slide3)
    
    add_title_text(slide3, "🚀 出道之路", Inches(0.5), font_size=40)
    
    timeline_points = [
        "2012年：参加《向上吧！少年》，进入全国200强（出道要从娃娃抓起）",
        "2015年：参加中韩节目《星动亚洲》，开始闪闪发光 ✨",
        "2016年：作为SWIN组合成员正式出道（没听过？正常，我也没听过）",
        "2018年：《偶像练习生》第一名！组建NINE PERCENT出道 🎉",
        "2018年至今：发行多张专辑，参演综艺，成为顶流偶像",
    ]
    add_bullet_points(slide3, timeline_points, Inches(1.5), font_size=26)
    
    add_emoji_text(slide3, "🔥", "从练习生到顶流，只用了不到6年！", Inches(6.3), font_size=20)
    
    # ========== 第4页：视频播放页 ==========
    slide4 = prs.slides.add_slide(blank_layout)
    create_gradient_slide(slide4)
    
    add_title_text(slide4, "🎬 传说中的「鸡你太美」", Inches(0.8), font_size=44)
    
    # 视频说明区域
    video_info_box = slide4.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3.5))
    tf = video_info_box.text_frame
    tf.word_wrap = True
    
    lines = [
        ("🎵 视频：《只因你太美》原版完整版现场", 32, RGBColor(255, 200, 100)),
        ("", 16, RGBColor(200, 200, 200)),
        ("📺 时长：1分钟 | 画质：720P", 24, RGBColor(200, 200, 200)),
        ("", 16, RGBColor(200, 200, 200)),
        ("💡 小贴士：请在演示模式下播放视频", 22, RGBColor(150, 200, 255)),
        ("    视频文件：caixukun_jntm.mp4", 20, RGBColor(180, 180, 180)),
        ("", 16, RGBColor(200, 200, 200)),
        ("📁 请确保视频文件与PPT在同一目录下", 20, RGBColor(255, 180, 180)),
    ]
    
    for i, (text, size, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "PingFang SC"
        p.space_after = Pt(8)
    
    # 添加播放提示框
    play_box = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4.5), Inches(5.8),
        Inches(4.33), Inches(1)
    )
    play_box.fill.solid()
    play_box.fill.fore_color.rgb = RGBColor(200, 50, 50)
    play_box.line.fill.background()
    
    play_text = slide4.shapes.add_textbox(Inches(4.5), Inches(5.95), Inches(4.33), Inches(0.8))
    tf = play_text.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "▶️ 点击播放视频"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)
    run.font.name = "PingFang SC"
    
    # ========== 第5页：代表作与名场面 ==========
    slide5 = prs.slides.add_slide(blank_layout)
    create_gradient_slide(slide5)
    
    add_title_text(slide5, "🎵 代表作 & 梗解读", Inches(0.3), font_size=40)
    
    # 左侧：音乐作品
    left_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.5), Inches(3))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "🎤 音乐作品"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(100, 200, 255)
    run.font.name = "PingFang SC"
    
    songs = ["《Wait Wait Wait》", "《情人》", "《YOUNG》", "《迷》", "《Hard To Get》"]
    for song in songs:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = f"  • {song}"
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(220, 220, 220)
        run.font.name = "PingFang SC"
        p.space_after = Pt(6)
    
    # 右侧：梗解读
    right_box = slide5.shapes.add_textbox(Inches(6.5), Inches(1.2), Inches(6), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "🐔 「鸡你太美」梗解读"
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 180, 100)
    run.font.name = "PingFang SC"
    
    meme_text = [
        "",
        "来源：《偶像练习生》表演曲《Rulez》",
        "",
        "事件：一段篮球舞被网友疯狂恶搞",
        "",
        "结果：成为中文互联网最火表情包之一",
        "",
        "态度：本人在综艺中多次自嘲，",
        "展示了强大的心理素质 💪",
    ]
    for line in meme_text:
        p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(20)
        run.font.color.rgb = RGBColor(220, 220, 220)
        run.font.name = "PingFang SC"
        p.space_after = Pt(4)
    
    add_emoji_text(slide5, "🐔", "「只因你太美~」—— 互联网永远的神", Inches(6.5), font_size=18)
    
    # ========== 第6页：总结 ==========
    slide6 = prs.slides.add_slide(blank_layout)
    create_gradient_slide(slide6)
    
    add_title_text(slide6, "📝 一句话总结", Inches(1), font_size=40)
    
    add_title_text(slide6, 
        "「他是一个被鬼畜成就的偶像，\n也是一个用作品证明自己的艺人」",
        Inches(2.2), font_size=32, bold=False, color=RGBColor(255, 215, 0))
    
    summary_points = [
        "🎵 音乐实力：确实会唱歌、会跳舞、会创作",
        "🏀 篮球技术：......咱就不评价了",
        "💪 心理素质：面对全网玩梗依然微笑面对",
        "📈 商业价值：顶流就是顶流，带货能力杠杠的",
    ]
    add_bullet_points(slide6, summary_points, Inches(4), left=Inches(2), font_size=24)
    
    add_emoji_text(slide6, "🎤", "谢谢大家！记得给我点赞~", Inches(6.5), font_size=24)
    
    # 保存PPT
    output_path = "/Users/lillianliao/notion_rag/蔡徐坤介绍_含视频版.pptx"
    prs.save(output_path)
    print(f"✅ PPT已生成: {output_path}")
    print(f"📁 请确保视频文件 caixukun_jntm.mp4 与PPT放在同一目录下")
    print(f"💡 提示：在PowerPoint中手动插入视频：插入 -> 视频 -> 选择 caixukun_jntm.mp4")
    return output_path

if __name__ == "__main__":
    create_caixukun_ppt_with_video()
